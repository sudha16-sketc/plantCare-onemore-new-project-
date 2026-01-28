"""
PPT Visual Guide Generator Service

Generates a PowerPoint (.pptx) plant care guide
using structured data from GeminiResponse.

This replaces NotebookLM during development and can be
swapped out for actual NotebookLM API integration later.
"""

import os
import asyncio
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt

from app.schemas.plant import GeminiResponse, NotebookLMResponse


class PPTService:
    """
    Generates PPT-based visual guides as a replacement
    for NotebookLM during the hackathon.
    """

    def __init__(self):
        self.output_dir = "generated_files"
        os.makedirs(self.output_dir, exist_ok=True)

    async def generate(
        self,
        plant_care_data: GeminiResponse,
        plant_name: str
    ) -> NotebookLMResponse:
        """
        Generate a PowerPoint presentation from plant care data.
        """

        filename = f"{plant_name.lower().replace(' ', '_')}_care_guide.pptx"
        file_path = os.path.join(self.output_dir, filename)

        prs = Presentation()

        # ---------- Slide 1: Title ----------
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = f"{plant_name} Care Guide"
        slide.placeholders[1].text = (
            f"Difficulty: {plant_care_data.plant_overview.difficulty_level}"
        )

        # ---------- Slide 2: Plant Overview ----------
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Plant Overview"
        slide.placeholders[1].text = plant_care_data.plant_overview.description

        # ---------- Slide 3: Ideal Conditions ----------
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Ideal Growing Conditions"

        conditions = plant_care_data.plant_overview.ideal_conditions
        # Handle dict access - ideal_conditions is a dict, not an object
        temp = conditions.get('temperature', 'Not specified') if isinstance(conditions, dict) else getattr(conditions, 'temperature', 'Not specified')
        humidity = conditions.get('humidity', 'Not specified') if isinstance(conditions, dict) else getattr(conditions, 'humidity', 'Not specified')
        sunlight = conditions.get('sunlight', 'Not specified') if isinstance(conditions, dict) else getattr(conditions, 'sunlight', 'Not specified')
        soil_ph = conditions.get('soil_ph', 'Not specified') if isinstance(conditions, dict) else getattr(conditions, 'soil_ph', 'Not specified')
        
        slide.placeholders[1].text = (
            f"Temperature: {temp}\n"
            f"Humidity: {humidity}\n"
            f"Sunlight: {sunlight}\n"
            f"Soil pH: {soil_ph}"
        )

        # ---------- Growth Stages ----------
        for stage in plant_care_data.growth_stages:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = stage.stage_name
            slide.placeholders[1].text = (
                f"Duration: {stage.duration}\n\n"
                f"Care:\n{stage.care_instructions}\n\n"
                f"Indicators:\n- " + "\n- ".join(stage.key_indicators)
            )

        # ---------- Daily Care ----------
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Daily Care Routine"

        daily = plant_care_data.daily_care
        slide.placeholders[1].text = (
            "Morning:\n- " + "\n- ".join(daily.morning_routine) + "\n\n"
            "Afternoon:\n- " + "\n- ".join(daily.afternoon_routine) + "\n\n"
            "Evening:\n- " + "\n- ".join(daily.evening_routine)
        )

        # ---------- Common Problems ----------
        for problem in plant_care_data.common_problems:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = problem.problem
            slide.placeholders[1].text = (
                "Symptoms:\n- " + "\n- ".join(problem.symptoms) + "\n\n"
                f"Solution:\n{problem.solution}\n\n"
                f"Prevention:\n{problem.prevention}"
            )

        # ---------- Tips ----------
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Additional Tips"
        slide.placeholders[1].text = "- " + "\n- ".join(
            plant_care_data.additional_tips
        )

        # Save PPT
        prs.save(file_path)

        return NotebookLMResponse(
            status="success",
            file_url=f"/files/{filename}",
            file_type="pptx",
            message="PPT visual guide generated successfully"
        )


# Singleton
ppt_service = PPTService()
